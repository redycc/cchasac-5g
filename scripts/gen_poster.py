"""
A0 conference poster — H-RB / C-HASAC DL Final Project
Story arc (Option A):
  HASAC (JMLR 2024) as base  →  apply to 5G PF coordination
  →  discover structural limitation  →  H-RB extends HASAC with discrete manager

Layout: Title → 2-col (Intro+HASAC baseline | H-RB Methodology) → Full-width Results → Conclusion
Design: Times New Roman, FACADE spec (Title ~80pt, body 40pt, caption 28pt italic)
Logos: transparent PNG (NYCU left, WinLab right)
"""
import os
from pptx import Presentation
from pptx.util import Mm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ───────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1A, 0x35, 0x6E)
BLUE   = RGBColor(0x1A, 0x5C, 0xB8)
TEAL   = RGBColor(0x00, 0x73, 0x85)
GREEN  = RGBColor(0x17, 0x73, 0x3A)
RED    = RGBColor(0xAA, 0x28, 0x20)
GOLD   = RGBColor(0xB8, 0x86, 0x00)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BLACK  = RGBColor(0x18, 0x18, 0x18)
GRAY   = RGBColor(0x55, 0x55, 0x55)
LGRAY  = RGBColor(0xF4, 0xF6, 0xFA)
SLBLUE = RGBColor(0xE6, 0xEF, 0xFA)
LGREEN = RGBColor(0xE2, 0xF5, 0xE9)
LYELLO = RGBColor(0xFD, 0xF5, 0xDC)
LPINK  = RGBColor(0xFC, 0xEC, 0xEB)

F = "Times New Roman"

# ── Slide ─────────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Mm(841)
prs.slide_height = Mm(1189)
W = prs.slide_width
H = prs.slide_height
slide = prs.slides.add_slide(prs.slide_layouts[6])

def bg_white():
    s = slide.shapes.add_shape(1, 0, 0, W, H)
    s.fill.solid(); s.fill.fore_color.rgb = WHITE
    s.line.fill.background()
bg_white()

# ── Primitives ────────────────────────────────────────────────────────────────
def box(x, y, w, h, fill=None, line=None, lw=0.5):
    s = slide.shapes.add_shape(1, x, y, w, h)
    if fill: s.fill.solid(); s.fill.fore_color.rgb = fill
    else: s.fill.background()
    if line: s.line.color.rgb = line; s.line.width = Pt(lw)
    else: s.line.fill.background()
    return s

def pic(path, x, y, w=None, h=None):
    if not os.path.exists(path): return
    from PIL import Image as PI
    img = PI.open(path); iw, ih = img.size
    if w and not h: h = int(w * ih / iw)
    elif h and not w: w = int(h * iw / ih)
    slide.shapes.add_picture(path, x, y, width=w, height=h)

def tb(x, y, w, h, rows, wrap=True):
    t = slide.shapes.add_textbox(x, y, w, h)
    tf = t.text_frame; tf.word_wrap = wrap
    first = True
    for r in rows:
        txt = r.get("t",""); sz = r.get("sz",40); b = r.get("b",False)
        i = r.get("i",False); c = r.get("c",BLACK); a = r.get("a",PP_ALIGN.LEFT)
        sp = r.get("sp",0)
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.alignment = a
        if sp: p.space_before = Pt(sp)
        run = p.add_run(); run.text = txt
        run.font.name = F; run.font.size = Pt(sz)
        run.font.bold = b; run.font.italic = i
        run.font.color.rgb = c
    return t

def R(t="", sz=40, b=False, i=False, c=BLACK, a=PP_ALIGN.LEFT, sp=0):
    return dict(t=t, sz=sz, b=b, i=i, c=c, a=a, sp=sp)

def sec_band(x, y, w, label, color=NAVY):
    bh = Mm(20)
    box(x, y, w, bh, fill=color)
    tb(x+Mm(6), y+Mm(2), w-Mm(12), bh,
       [R(label, sz=54, b=True, c=WHITE)])
    return y + bh

def callout(x, y, w, h, color, lines):
    """Coloured left-border callout box."""
    box(x, y, Mm(5), h, fill=color)
    box(x+Mm(5), y, w-Mm(5), h, fill=RGBColor(
        min(color.rgb >> 16, 0xFF),
        min((color.rgb >> 8) & 0xFF, 0xFF),
        min(color.rgb & 0xFF, 0xFF),
    ))
    # lighter fill
    r2 = min(((color.rgb >> 16) & 0xFF) + 90, 255)
    g2 = min(((color.rgb >> 8)  & 0xFF) + 90, 255)
    b2 = min( (color.rgb & 0xFF)         + 90, 255)
    box(x+Mm(5), y, w-Mm(5), h, fill=RGBColor(r2, g2, b2))
    tb(x+Mm(9), y+Mm(4), w-Mm(13), h-Mm(6), lines)

# ── Layout dims ───────────────────────────────────────────────────────────────
M   = Mm(22)
GH  = Mm(14)
GV  = Mm(16)
TH  = Mm(112)
CW2 = (W - 2*M - GH) // 2

# ═══════════════════════════════════════════════════════════════════════════════
# TITLE BAR
# ═══════════════════════════════════════════════════════════════════════════════
box(0, 0, W, TH, fill=NAVY)

LH = Mm(72)
LW_NYCU = int(LH * 805/806)
LW_WIN  = int(LH * 164/225)
pic("logo/guideline_logo0.49478bb3.png", M, Mm(20), h=LH)
pic("logo/WinLogo_transparent.png", W-M-LW_WIN, Mm(24), h=LH)

tx = M + LW_NYCU + Mm(18)
tw = W - 2*M - LW_NYCU - LW_WIN - Mm(36)
tb(tx, Mm(14), tw, Mm(52), [
    R("Extending HASAC with Hierarchical Discrete Management", sz=72, b=True, c=WHITE, a=PP_ALIGN.CENTER),
    R("for 5G Multi-Cell Resource Allocation", sz=72, b=True, c=WHITE, a=PP_ALIGN.CENTER),
])
tb(tx, Mm(67), tw, Mm(22), [
    R("When the Action Space Is the Bottleneck: A Structural Diagnosis and Fix",
      sz=40, i=True, c=RGBColor(0xBB,0xD4,0xFF), a=PP_ALIGN.CENTER),
])
tb(tx, Mm(91), tw, Mm(16), [
    R("DL Final Project  ·  WINLAB, National Yang Ming Chiao Tung University  ·  2026",
      sz=33, c=RGBColor(0x88,0xAA,0xDD), a=PP_ALIGN.CENTER),
])

# ═══════════════════════════════════════════════════════════════════════════════
# ROW 1  — left: Introduction & HASAC baseline   right: H-RB Methodology
# ═══════════════════════════════════════════════════════════════════════════════
R1_Y = TH + GV
R1_H = Mm(444)

# ─── LEFT: Introduction ──────────────────────────────────────────────────────
x1 = M;  y = R1_Y
box(x1, y, CW2, R1_H, fill=LGRAY)
y = sec_band(x1, y, CW2, "① Background & Motivation", color=NAVY)
y += Mm(10)

# 1-a  Problem
tb(x1+Mm(8), y, CW2-Mm(16), Mm(24),[
    R("The 5G Coordination Problem", sz=46, b=True, c=NAVY),
])
y += Mm(26)
tb(x1+Mm(8), y, CW2-Mm(16), Mm(60), [
    R("Multiple base stations (BS) share spectrum. Simultaneous", sz=40),
    R("transmissions cause inter-cell interference. The goal:", sz=40),
])
y += Mm(44)
box(x1+Mm(16), y, CW2-Mm(32), Mm(22), fill=SLBLUE, line=BLUE, lw=1.5)
tb(x1+Mm(22), y+Mm(3), CW2-Mm(44), Mm(16), [
    R("Maximise PF-utility  U = Σ_u log(R̄_u)    [floor ≈ −6.3,  oracle ≈ +23.7]",
      sz=38, b=True, c=NAVY, a=PP_ALIGN.CENTER),
])
y += Mm(32)
tb(x1+Mm(8), y, CW2-Mm(16), Mm(50), [
    R("Setting:  3 BS · 10 UE · 4 RB  |  each BS sees local obs only.", sz=40),
    R("Dynamic env: UE random walk + per-UE queue buffer (HOL latency).", sz=40),
    R("Evaluation also measures Goodput and P99 HOL latency.", sz=40),
])
y += Mm(58)

# 1-b  Our starting point: HASAC
tb(x1+Mm(8), y, CW2-Mm(16), Mm(24), [
    R("Starting Point: HASAC  [HARL, JMLR 2024]", sz=46, b=True, c=BLUE),
])
y += Mm(26)
box(x1+Mm(8), y, CW2-Mm(16), Mm(66), fill=RGBColor(0xE6,0xEF,0xFA), line=BLUE, lw=1.5)
tb(x1+Mm(14), y+Mm(5), CW2-Mm(28), Mm(58), [
    R("Heterogeneous-Agent SAC — key properties:", sz=40, b=True, c=BLUE),
    R("  • Sequential update scheme (monotonic improvement guarantee)", sz=38),
    R("  • CTDE: centralized training, decentralized execution", sz=38),
    R("  • No parameter sharing — supports heterogeneous agents", sz=38),
    R("  • Off-policy SAC with shared twin-Q critic", sz=38),
    R("We adopt HASAC workers + shared twin-Q as our base agent.", sz=38, b=True, c=NAVY),
])
y += Mm(74)

# 1-c  What we found: HASAC hits a wall
tb(x1+Mm(8), y, CW2-Mm(16), Mm(24), [
    R("Finding: HASAC Hits a Hard Ceiling", sz=46, b=True, c=RED),
])
y += Mm(26)

# comparison callout
box(x1+Mm(8), y, CW2-Mm(16), Mm(60), fill=RGBColor(0xFC,0xEC,0xEB), line=RED, lw=1.5)
tb(x1+Mm(14), y+Mm(5), CW2-Mm(28), Mm(52), [
    R("29 variants of HASAC (cc-HASAC v1–v24 + goodput v1–v5):", sz=40, b=True),
    R("  Max PF-U ever = −4.26    Max goodput = 26.15", sz=40, c=RED, b=True),
    R("", sz=6),
    R("Fixed-power constant  (zero RL):  PF-U = −4.17", sz=38, i=True),
    R("→ 200k-step HASAC training ≈ a constant baseline.", sz=40, b=True, c=RED),
])
y += Mm(68)

# 1-d  Root-cause
tb(x1+Mm(8), y, CW2-Mm(16), Mm(24), [
    R("Root Cause: Structural Mismatch", sz=46, b=True, c=TEAL),
])
y += Mm(26)
tb(x1+Mm(8), y, CW2-Mm(16), Mm(84), [
    R("① SAC tanh saturation  →  −165 crash", sz=42, b=True),
    R("   mu → −∞ → power_frac → 0 → rates ≈ 0.", sz=38, i=True, c=GRAY),
    R("   Fix: mu = 5·tanh(mu_raw).  Confirmed via pwr log.", sz=38, c=TEAL),
    R("", sz=6),
    R("② Continuous per-UE action space is wrong for PF", sz=42, b=True),
    R("   +23.7 oracle gain comes from cooperative power back-off.", sz=38, i=True, c=GRAY),
    R("   SAC entropy destroys any coordination once found.", sz=38, i=True, c=GRAY),
    R("", sz=6),
    R("③ Airtight proof (centralized full-CSI SAC): −5.7", sz=42, b=True),
    R("   Even with full information, SAC cannot learn PF coordination.", sz=38, i=True, c=GRAY),
    R("→ Not a tuning problem.  An action-space structure problem.",
      sz=42, b=True, c=NAVY),
])

# ─── RIGHT: Methodology ──────────────────────────────────────────────────────
x2 = M + CW2 + GH;  y = R1_Y
box(x2, y, CW2, R1_H, fill=LGRAY)
y = sec_band(x2, y, CW2, "② Methodology: Extending HASAC with H-RB", color=TEAL)
y += Mm(10)

# 2-a  Insight
tb(x2+Mm(8), y, CW2-Mm(16), Mm(24), [
    R("Structural Insight from Oracle Analysis", sz=46, b=True, c=TEAL),
])
y += Mm(26)
box(x2+Mm(8), y, CW2-Mm(16), Mm(50), fill=RGBColor(0xE0,0xF4,0xF7), line=TEAL, lw=1.5)
tb(x2+Mm(14), y+Mm(5), CW2-Mm(28), Mm(42), [
    R("PF-WSR oracle (+23.7) wins by:", sz=42, b=True, c=TEAL),
    R("  ① Grid-searching each BS's discrete total power level", sz=40),
    R("       (cooperative back-off across cells)", sz=38, i=True, c=GRAY),
    R("  ② Analytical weighted water-filling within each budget", sz=40),
    R("→ Coordinaton is fully captured by one discrete choice per BS.",
      sz=40, b=True, c=NAVY),
])
y += Mm(58)

# 2-b  H-RB design
tb(x2+Mm(8), y, CW2-Mm(16), Mm(24), [
    R("H-RB: Hierarchical Extension of HASAC", sz=46, b=True, c=NAVY),
])
y += Mm(26)

AX = x2+Mm(12); AW = CW2-Mm(24)

# Manager box
box(AX, y, AW, Mm(82), fill=NAVY)
tb(AX+Mm(6), y+Mm(6), AW-Mm(12), Mm(72), [
    R("MANAGER  ·  new discrete coordinator  (every K = 10 steps)",
      sz=41, b=True, c=WHITE),
    R("", sz=5),
    R("obs:  global KPM — SINR / load / goodput / buf / HOL  [27-dim]",
      sz=36, c=RGBColor(0xBB,0xD6,0xFF)),
    R("act:  categorical per RB  →  assignment[rb] ∈ {0,1,2}",
      sz=36, c=RGBColor(0xBB,0xD6,0xFF)),
    R("algo: factored discrete SAC  (per-RB head + twin-Q + auto-α)",
      sz=36, c=RGBColor(0xBB,0xD6,0xFF)),
    R("role: slow-timescale RB partition  (xApp via O-RAN E2 interface)",
      sz=33, i=True, c=GOLD),
])
y += Mm(82)

# Arrow
tb(AX, y+Mm(2), AW, Mm(13), [
    R("↓  Hard RB ownership mask  —  orthogonality enforced by construction",
      sz=35, b=True, c=NAVY, a=PP_ALIGN.CENTER),
])
y += Mm(17)

# Workers box
box(AX, y, AW, Mm(82), fill=BLUE)
tb(AX+Mm(6), y+Mm(6), AW-Mm(12), Mm(72), [
    R("WORKERS  ·  HASAC agents  (every step)",
      sz=41, b=True, c=WHITE),
    R("", sz=5),
    R("obs:  local KPM(9) + RB-ownership mask(4) + agent-id(3)  [16-dim]",
      sz=36, c=RGBColor(0xCC,0xE5,0xFF)),
    R("act:  continuous per-RB power;  non-owned RBs forced to zero",
      sz=36, c=RGBColor(0xCC,0xE5,0xFF)),
    R("algo: continuous SAC  +  shared twin-Q  (CTDE)  ← HASAC core",
      sz=36, b=True, c=GOLD),
    R("role: fast-timescale power control  (gNB)",
      sz=33, i=True, c=RGBColor(0xCC,0xE5,0xFF)),
])
y += Mm(82)

tb(AX, y+Mm(2), AW, Mm(13), [
    R("↕  reward / observation", sz=35, b=True, c=NAVY, a=PP_ALIGN.CENTER),
])
y += Mm(17)

# Env box
box(AX, y, AW, Mm(52), fill=TEAL)
tb(AX+Mm(6), y+Mm(5), AW-Mm(12), Mm(44), [
    R("ENVIRONMENT  (cc_env_goodput_v2)", sz=41, b=True, c=WHITE),
    R("3 BS × 10 UE × 4 RB  ·  dynamic channels  ·  HOL queue buffer",
      sz=36, c=RGBColor(0xCC,0xF0,0xF5)),
    R("reward = log(1+thr) − β·(Q/Q_ref) − η·power",
      sz=34, i=True, c=RGBColor(0xCC,0xF0,0xF5)),
])
y += Mm(52) + Mm(14)

# 2-c  What changes vs vanilla HASAC
tb(x2+Mm(8), y, CW2-Mm(16), Mm(24), [
    R("H-RB vs Vanilla HASAC", sz=46, b=True, c=NAVY),
])
y += Mm(26)

diff_rows = [
    ("",            "Vanilla HASAC",         "H-RB (ours)",      True),
    ("Manager",     "none",                  "Discrete SAC ★",   False),
    ("Worker algo", "Continuous SAC",        "Continuous SAC ✓", False),
    ("Worker act",  "full per-RB power",     "power on owned RBs only", False),
    ("Coord. mech", "broadcast z / none",    "hard RB partition",False),
    ("CTDE",        "✓",                     "✓",                False),
    ("Seq. update", "✓",                     "✓",                False),
]
cols_d = [CW2*36//100, CW2*30//100, CW2*28//100]
rx = x2+Mm(8); ry = y
for i, (lbl, v1, v2, hdr) in enumerate(diff_rows):
    bg = NAVY if hdr else (SLBLUE if i%2==0 else WHITE)
    fc = WHITE if hdr else BLACK
    xx = rx
    for wi, txt in zip(cols_d, [lbl, v1, v2]):
        box(xx, ry, wi, Mm(13), fill=bg, line=RGBColor(0xAA,0xBB,0xCC))
        tb(xx+Mm(2), ry+Mm(1.5), wi-Mm(4), Mm(12),
           [R(txt, sz=30, b=hdr, c=GREEN if "★" in txt else fc)])
        xx += wi
    ry += Mm(13)

# ═══════════════════════════════════════════════════════════════════════════════
# ROW 2  — Key Results  (full width, 3 equal columns)
# ═══════════════════════════════════════════════════════════════════════════════
R2_Y = R1_Y + R1_H + GV
R2_H = Mm(356)
GH3  = Mm(12)
CW3  = (W - 2*M - 2*GH3) // 3

box(M, R2_Y, W-2*M, R2_H, fill=LGRAY)
ys = sec_band(M, R2_Y, W-2*M, "③ Key Results", color=GREEN)

# ── 3a: H-RB Goodput ─────────────────────────────────────────────────────────
xa = M; ya = ys+Mm(10)
tb(xa+Mm(8), ya, CW3-Mm(16), Mm(24), [
    R("H-RB: Goodput & P99 Latency", sz=44, b=True, c=GREEN),
])
ya += Mm(26)
box(xa+Mm(8), ya, CW3-Mm(16), Mm(44), fill=LGREEN, line=GREEN, lw=2)
tb(xa+Mm(14), ya+Mm(5), CW3-Mm(28), Mm(36), [
    R("vs flat HASAC best:    Goodput  +14.5%", sz=41, b=True, c=GREEN, a=PP_ALIGN.CENTER),
    R("26.15  →  29.96 bits/step      P99  −49%", sz=41, b=True, c=GREEN, a=PP_ALIGN.CENTER),
])
ya += Mm(50)

hrb = [
    ("Strategy",                  "Goodput", "P99",   True,  NAVY),
    ("Full-power  (no coord)",     "20.52",  "99.5",  False, BLACK),
    ("HASAC goodput v5  (flat)",   "26.15",  "99.5",  False, RED),
    ("Random manager",             "29.27",  "66.6",  False, BLACK),
    ("Freq-reuse oracle",          "29.96",  "56.1",  False, GRAY),
    ("Static partition",           "29.96",  "52.6",  False, GRAY),
    ("★  H-RB learned",           "29.96",  "50.5",  True,  GREEN),
]
cols_h = [CW3*56//100, CW3*22//100, CW3*16//100]
rx = xa+Mm(8); ry = ya
for i,(lbl,gp,p99,b,c) in enumerate(hrb):
    bg = SLBLUE if i%2==0 else WHITE
    if "H-RB learned" in lbl: bg = LGREEN
    if "HASAC" in lbl and "flat" in lbl: bg = RGBColor(0xFC,0xEC,0xEB)
    xx = rx
    for wi,txt in zip(cols_h,[lbl,gp,p99]):
        box(xx,ry,wi,Mm(14),fill=bg,line=RGBColor(0xBB,0xCC,0xBB))
        tb(xx+Mm(2),ry+Mm(1.5),wi-Mm(4),Mm(13),[R(txt,sz=30,b=b,c=c)])
        xx += wi
    ry += Mm(14)
ya = ry+Mm(10)
tb(xa+Mm(8), ya, CW3-Mm(16), Mm(52), [
    R("Dynamic RB reallocation (learned manager) achieves", sz=38),
    R("P99 = 50.5 — lower than any fixed partition.", sz=38),
    R("This is the sequential decision advantage HASAC", sz=38),
    R("workers provide once given the right action space.", sz=38, b=True, c=NAVY),
])

# ── 3b: PF-Utility ───────────────────────────────────────────────────────────
xb = M+CW3+GH3; yb = ys+Mm(10)
tb(xb+Mm(8), yb, CW3-Mm(16), Mm(24), [
    R("PF-Utility: HASAC vs H-RB Structure", sz=44, b=True, c=GREEN),
])
yb += Mm(26)
box(xb+Mm(8), yb, CW3-Mm(16), Mm(44), fill=LGREEN, line=GREEN, lw=2)
tb(xb+Mm(14), yb+Mm(5), CW3-Mm(28), Mm(36), [
    R("HASAC (flat):  −4.26    →    H-RB random:  −0.71", sz=41, b=True, c=GREEN, a=PP_ALIGN.CENTER),
    R("Structure alone  =  +3.55 PF-U  (zero learning needed)", sz=39, b=True, c=NAVY, a=PP_ALIGN.CENTER),
])
yb += Mm(50)

pf = [
    ("Strategy",                      "PF-U",  True,  NAVY),
    ("Equal-power floor",              "−6.27", False, BLACK),
    ("HASAC z1 (best flat RL, 200k)",  "−4.26", False, RED),
    ("Centralized full-CSI SAC",       "−5.71", False, RED),
    ("Fixed power 0.75 (zero RL)",     "−4.17", False, GRAY),
    ("★  H-RB random manager",        "−0.71", True,  GREEN),
    ("H-RB oracle (reachable ceiling)","  +23.7", True, GREEN),
]
cols_p = [CW3*68//100, CW3*26//100]
rx2 = xb+Mm(8); ry2 = yb
for i,(lbl,val,b,c) in enumerate(pf):
    bg = SLBLUE if i%2==0 else WHITE
    if val in ("+23.7","  +23.7","−0.71"): bg = LGREEN
    if "HASAC z1" in lbl or "Centralized" in lbl: bg = RGBColor(0xFC,0xEC,0xEB)
    box(rx2,ry2,cols_p[0],Mm(14),fill=bg,line=RGBColor(0xBB,0xCC,0xBB))
    box(rx2+cols_p[0],ry2,cols_p[1],Mm(14),fill=bg,line=RGBColor(0xBB,0xCC,0xBB))
    tb(rx2+Mm(2),ry2+Mm(1.5),cols_p[0]-Mm(4),Mm(13),[R(lbl,sz=30,b=b,c=c)])
    tb(rx2+cols_p[0]+Mm(2),ry2+Mm(1.5),cols_p[1]-Mm(4),Mm(13),[R(val,sz=30,b=b,c=c)])
    ry2 += Mm(14)
yb = ry2+Mm(10)
tb(xb+Mm(8), yb, CW3-Mm(16), Mm(52), [
    R("Oracle bracket is [−0.71, +23.7] — both are reachable", sz=38),
    R("within the H-RB structure.  The learned manager", sz=38),
    R("(train-acc 0.80) overfits to zero-power combos; excluding", sz=38),
    R("level-0 is the immediate next fix.", sz=38, b=True, c=NAVY),
])

# ── 3c: cc-HASAC Sum-Rate ablation ───────────────────────────────────────────
xc = M+2*(CW3+GH3); yc = ys+Mm(10)
tb(xc+Mm(8), yc, CW3-Mm(16), Mm(24), [
    R("cc-HASAC: z Contribution Ablation", sz=44, b=True, c=GREEN),
])
yc += Mm(26)
box(xc+Mm(8), yc, CW3-Mm(16), Mm(44), fill=LYELLO, line=GOLD, lw=2)
tb(xc+Mm(14), yc+Mm(5), CW3-Mm(28), Mm(36), [
    R("Best (v22 TD3+BC):  49.35 bps/Hz", sz=41, b=True, c=GOLD, a=PP_ALIGN.CENTER),
    R("Ind-SAC A (no z, ref):  28.10  (+75.6%)", sz=39, b=True, c=NAVY, a=PP_ALIGN.CENTER),
])
yc += Mm(50)

sr = [
    ("Method",                     "Sum-rate","z←0 Δ", True,  NAVY),
    ("Ind-SAC A  (no z, ref)",      "28.10",  "—",     False, GRAY),
    ("cc-HASAC v6  (BC pretrain)",  "32.97",  "+6.14", False, BLACK),
    ("cc-HASAC v13  (best static)", "34.21",  "+4.35", True,  GREEN),
    ("cc-HASAC v17  (per-BS z)",    "39.98†", "+12.43",False, BLACK),
    ("★ cc-HASAC v22  (TD3+BC)",   "49.35†", "+21.61",True,  GOLD),
]
cols_s = [CW3*54//100, CW3*22//100, CW3*18//100]
rx3 = xc+Mm(8); ry3 = yc
for i,row in enumerate(sr):
    lbl,sv,dz,b,c = row
    bg = SLBLUE if i%2==0 else WHITE
    if "v22" in lbl: bg = LYELLO
    if "v13" in lbl: bg = LGREEN
    xx = rx3
    for wi,txt in zip(cols_s,[lbl,sv,dz]):
        box(xx,ry3,wi,Mm(14),fill=bg,line=RGBColor(0xBB,0xCC,0xDD))
        tb(xx+Mm(2),ry3+Mm(1.5),wi-Mm(4),Mm(13),[R(txt,sz=30,b=b,c=c)])
        xx += wi
    ry3 += Mm(14)
yc = ry3+Mm(8)
tb(xc+Mm(8), yc, CW3-Mm(16), Mm(18), [
    R("† Dynamic goodput env (not static sum-rate).", sz=30, i=True, c=GRAY),
])
yc += Mm(20)
tb(xc+Mm(8), yc, CW3-Mm(16), Mm(110), [
    R("Key design findings for HASAC workers:", sz=40, b=True, c=NAVY),
    R("", sz=6),
    R("Encoder freeze ablation:", sz=38, b=True),
    R("  10k  freeze  →  z←0 Δ = −6.1  (z harmful)", sz=36, c=RED),
    R("  100k freeze, lr 1e-4  →  Δ = +1.8  (weak)", sz=36),
    R("  100k freeze, lr 1e-5  →  Δ = +4.4  ★", sz=36, b=True, c=GREEN),
    R("", sz=6),
    R("Per-BS personalised z  (v17):", sz=38, b=True),
    R("  z←0 Δ = +12.43  ← strongest z contribution across all", sz=36, b=True, c=GREEN),
    R("  3-token broadcast → attention degenerates to mean-pool.", sz=36, c=GRAY),
    R("  Personalised z_i essential for N_BS = 3.", sz=36),
    R("", sz=6),
    R("RL vs BC tension:", sz=38, b=True),
    R("  BC policy alone reaches 48–51 bps/Hz.", sz=36),
    R("  RL phase crashes to 22–35 — Q overestimation.", sz=36, c=RED),
    R("  Best ckpt always from WARMUP (BC intact).", sz=36),
])

# ═══════════════════════════════════════════════════════════════════════════════
# ROW 3  — Conclusion / Contributions  (full width)
# ═══════════════════════════════════════════════════════════════════════════════
R3_Y = R2_Y + R2_H + GV
R3_H = Mm(186)
box(M, R3_Y, W-2*M, R3_H, fill=LGRAY)
ys3 = sec_band(M, R3_Y, W-2*M, "④ Conclusions & Contributions", color=NAVY)

CW4 = (W - 2*M - 3*GH3) // 4
ys3 += Mm(10)

cards = [
    ("①  Diagnosing\nHASAC's Limits",
     ["SAC tanh saturation → −165 crash",
      "(mu → −∞ → power → 0)",
      "Fix: mu-bound trick; confirmed",
      "via pwr diagnostic log.",
      "Continuous action space cannot",
      "express cooperative back-off."],
     TEAL, RGBColor(0xE0,0xF4,0xF7)),
    ("②  Negative Result\non HASAC",
     ["29 variants, all cap at PF-U ≈ −4.",
      "Equal to a constant baseline.",
      "Centralized full-CSI SAC: −5.7.",
      "Not a tuning problem —",
      "structural mismatch between",
      "SAC action space & PF objective."],
     RED, RGBColor(0xFC,0xEC,0xEB)),
    ("③  H-RB: Structural\nExtension of HASAC",
     ["Discrete manager + HASAC workers.",
      "Goodput +14.5%  (26.15→29.96)",
      "P99  −49%  (99.5→50.5 slots)",
      "PF-U: −4.26 → −0.71  (random!)",
      "Oracle bracket [−0.71, +23.7]",
      "reachable within structure."],
     GREEN, LGREEN),
    ("④  Design Lessons\nfor O-RAN xApp",
     ["Discrete coordinator (xApp, slow)",
      "handles inter-cell coordination.",
      "Analytic inner optimisation",
      "makes workers deployment-ready.",
      "Per-BS personalised z (enc lr 1e-5)",
      "is key for z contribution in HASAC."],
     BLUE, SLBLUE),
]

for k, (title, bullets, color, bg) in enumerate(cards):
    xk = M + k*(CW4+GH3)
    box(xk, ys3, CW4, R3_H-Mm(40), fill=bg, line=color, lw=2)
    box(xk, ys3, CW4, Mm(38), fill=color)
    lines = title.split("\n")
    tb(xk+Mm(5), ys3+Mm(4), CW4-Mm(10), Mm(30), [
        R(lines[0], sz=40, b=True, c=WHITE),
        R(lines[1] if len(lines)>1 else "", sz=40, b=True, c=WHITE),
    ])
    rows_b = [R("", sz=4)]
    for bul in bullets:
        rows_b.append(R("• " + bul, sz=36))
        rows_b.append(R("", sz=3))
    tb(xk+Mm(6), ys3+Mm(42), CW4-Mm(12), R3_H-Mm(82), rows_b)

# take-away strip
ty = R3_Y + R3_H - Mm(30)
box(M, ty, W-2*M, Mm(30), fill=NAVY)
tb(M+Mm(10), ty+Mm(6), W-2*M-Mm(20), Mm(20), [
    R("Take-away:  We adopt HASAC [HARL, JMLR 2024] as our base and show that its coordination bottleneck"
      " lies in action-space structure, not in algorithm tuning.  "
      "H-RB extends HASAC with a discrete hierarchical manager — structural gain alone"
      " (+3.55 PF-U, +14.5% goodput, −49% P99) exceeds 29 variants of reward and architecture search.",
      sz=35, b=True, c=WHITE, a=PP_ALIGN.CENTER),
])

# ── Footer ────────────────────────────────────────────────────────────────────
FY = H - Mm(22)
box(0, FY, W, Mm(22), fill=NAVY)
tb(M, FY+Mm(4), W-2*M, Mm(14), [
    R("Code: train_chasac_hrb.py · train_hier_rb.py · env_chasac.py · cc_env_goodput_v2.py"
      "  |  Base: HARL HASAC [JMLR 2024]  |  WINLAB, NYCU  ·  2026",
      sz=26, c=RGBColor(0x88,0xAA,0xDD), a=PP_ALIGN.CENTER),
])

# ── Save ─────────────────────────────────────────────────────────────────────
out = "/home/hyc1014/DL/FinalProject/poster_draft.pptx"
prs.save(out)
print(f"Saved → {out}")
