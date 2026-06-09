"""C-HASAC poster.pptx — A0 841×1189mm, Times New Roman.
Font spec:
  Title / Section Header : 84pt bold
  Author names           : 48pt
  Body text              : 42pt
  Detail / affiliation   : 28pt italic
"""
import os
from pptx import Presentation
from pptx.util import Mm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── palette ──────────────────────────────────────────────
C_PRIMARY = RGBColor(0x1a, 0x3a, 0x5c)
C_ACCENT  = RGBColor(0x2e, 0x86, 0xab)
C_SUCCESS = RGBColor(0x28, 0xa7, 0x45)
C_WARNING = RGBColor(0xdc, 0x35, 0x45)
C_ORANGE  = RGBColor(0xe6, 0x7e, 0x22)
C_PURPLE  = RGBColor(0x8e, 0x44, 0xad)
C_WHITE   = RGBColor(0xff, 0xff, 0xff)
C_TEXT    = RGBColor(0x21, 0x25, 0x29)
C_MUTED   = RGBColor(0x6c, 0x75, 0x7d)
C_CARD    = RGBColor(0xff, 0xff, 0xff)
C_BG      = RGBColor(0xf8, 0xf9, 0xfa)
C_BESTBG  = RGBColor(0xdf, 0xf0, 0xf8)
C_LIGHT   = RGBColor(0xf0, 0xf8, 0xff)
C_BORDER  = RGBColor(0xde, 0xe2, 0xe6)

HERE = os.path.dirname(os.path.abspath(__file__))
LOGO = os.path.join(HERE, '..', 'logo', 'WinLogo_transparent.png')

prs = Presentation()
prs.slide_width  = Mm(841)
prs.slide_height = Mm(1189)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

FONT = 'Times New Roman'

# ── primitive helpers ─────────────────────────────────────
def box(x, y, w, h, fill=None, line=None, lw=Pt(0.5)):
    shp = slide.shapes.add_shape(1, Mm(x), Mm(y), Mm(w), Mm(h))
    shp.fill.solid() if fill else shp.fill.background()
    if fill: shp.fill.fore_color.rgb = fill
    if line:
        shp.line.color.rgb = line; shp.line.width = lw
    else:
        shp.line.fill.background()
    return shp

def tb(x, y, w, h):
    """Create empty textbox, return text_frame."""
    txb = slide.shapes.add_textbox(Mm(x), Mm(y), Mm(w), Mm(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    return tf

def para(tf, text, size=Pt(42), bold=False, italic=False,
         color=C_TEXT, align=PP_ALIGN.LEFT, first=False, space_before=0):
    """Add (or reuse first) paragraph with a single run."""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.name  = FONT
    run.font.size  = size
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p

def single(x, y, w, h, text, size=Pt(42), bold=False, italic=False,
           color=C_TEXT, align=PP_ALIGN.LEFT):
    tf = tb(x, y, w, h)
    para(tf, text, size=size, bold=bold, italic=italic, color=color, align=align, first=True)

# ── layout ────────────────────────────────────────────────
M   = 18            # margin mm
GAP = 10            # column gap mm
W   = 841 - 2*M     # 805 mm usable width
COL = (W - GAP) / 2 # ~397.5 mm
LX  = M
RX  = M + COL + GAP

SHH = 42   # section header height mm  (84pt + padding)
SGAP = 8   # gap between sections mm

# ══════════════════════════════════════════════════════════
# HEADER  y=18, h=100mm
# ══════════════════════════════════════════════════════════
box(M, M, W, 100, fill=C_BG)

try:
    slide.shapes.add_picture(LOGO, Mm(LX+2), Mm(M+12), height=Mm(46))
except Exception:
    single(LX+2, M+14, 68, 24, 'WinLab', size=Pt(36), bold=True, color=C_PRIMARY)

# Title (84pt, 2 lines)
tf = tb(LX+78, M+2, W-162, 58)
para(tf, 'C-HASAC: Contextual Heterogeneous-Agent SAC',
     size=Pt(84), bold=True, color=C_PRIMARY, align=PP_ALIGN.CENTER, first=True)
para(tf, 'for Decentralized 5G Power Coordination via Learned Latent Context',
     size=Pt(60), bold=True, color=C_PRIMARY, align=PP_ALIGN.CENTER, space_before=4)

# Authors (48pt) — Yi-Chieh Hong underlined as presenter
txb_a = slide.shapes.add_textbox(Mm(LX+78), Mm(M+64), Mm(W-162), Mm(18))
txb_a.word_wrap = True
tf_a = txb_a.text_frame
tf_a.word_wrap = True
p_a = tf_a.paragraphs[0]
p_a.alignment = PP_ALIGN.CENTER
for i, (name, is_presenter) in enumerate([
    ('Yi-Chieh Hong', True),
    (',  Yen-Ting Kuo', False),
    (',  Wen-Ju Chiang', False),
]):
    run = p_a.add_run()
    run.text = name
    run.font.name = FONT
    run.font.size = Pt(48)
    run.font.bold = False
    run.font.color.rgb = C_TEXT
    if is_presenter:
        run.font.underline = True

# Affiliation (28pt italic)
single(LX+78, M+84, W-162, 14,
       'WinLab, National Yang Ming Chiao Tung University (NYCU)  ·  535518 Deep Learning Final Project, 2026',
       size=Pt(28), italic=True, color=C_MUTED, align=PP_ALIGN.CENTER)

# NYCU right label
tf_nycu = tb(M+W-82, M+8, 80, 46)
para(tf_nycu, 'National Yang Ming', size=Pt(30), bold=True, color=C_PRIMARY, align=PP_ALIGN.RIGHT, first=True)
para(tf_nycu, 'Chiao Tung University', size=Pt(30), bold=True, color=C_PRIMARY, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════
# BODY  starts at BY
# ══════════════════════════════════════════════════════════
BY = M + 100 + 8   # = 126mm

def section_hdr(x, y, w, title):
    box(x, y, w, SHH, fill=C_PRIMARY)
    single(x+4, y+4, w-8, SHH-8, title, size=Pt(84), bold=True, color=C_WHITE)

# ──────────────────────────────────────────────────────────
# LEFT COLUMN
# ──────────────────────────────────────────────────────────
y = BY

# ── BACKGROUND & MOTIVATION  ──────────────────────────────
BH = 255
section_hdr(LX, y, COL, 'Background & Motivation')
box(LX, y+SHH, COL, BH, fill=C_CARD, line=C_ACCENT, lw=Pt(0.4))

tf = tb(LX+5, y+SHH+6, COL-10, BH-12)
para(tf, '• 5G multi-cell networks: inter-cell interference degrades PF fairness.',
     size=Pt(42), color=C_TEXT, first=True)
para(tf, '• Goal: maximize Proportional-Fair utility across all UEs:',
     size=Pt(42), color=C_TEXT, space_before=6)
para(tf, '        U  =  Σu  log( R̄u + ε )',
     size=Pt(48), bold=True, color=C_PRIMARY, align=PP_ALIGN.CENTER, space_before=4)
para(tf, '• O-RAN constraint: each gNB observes only local UE metrics at runtime.',
     size=Pt(42), color=C_TEXT, space_before=6)
para(tf, '  No direct neighbor CSI exchange is permitted.',
     size=Pt(42), color=C_TEXT)
para(tf, '• RIC xApp computes latent context z from KPM measurements',
     size=Pt(42), color=C_TEXT, space_before=6)
para(tf, '  and pushes to gNBs via E2 interface — coordination without CSI leakage.',
     size=Pt(42), color=C_ACCENT)
para(tf, '── O-RAN: RIC ↔ E2 ↔ BS1 / BS2 / BS3  (KPM↑, z↓) ──',
     size=Pt(28), italic=True, color=C_MUTED, align=PP_ALIGN.CENTER, space_before=10)

y += SHH + BH + SGAP

# ── ARCHITECTURE  ─────────────────────────────────────────
AH = 350
section_hdr(LX, y, COL, 'Architecture')
box(LX, y+SHH, COL, AH, fill=C_CARD, line=C_ACCENT, lw=Pt(0.4))

# Three-layer subheading
single(LX+5, y+SHH+6, COL-10, 14,
       'Three-Layer Information Flow (strict separation)',
       size=Pt(42), bold=True, color=C_PRIMARY)

# Layer badges + text
layer_data = [
    ('A', C_ACCENT,  'BS-local obs',
     'per-UE {rate, PF weight, power}  →  Actor input (gNB-measurable in O-RAN)'),
    ('B', C_ORANGE,  'RIC KPM',
     '{load, throughput, P_BS, BS-distances} × 3  →  Encoder → z[16] → all Actors'),
    ('C', C_PURPLE,  'Privileged',
     'full CSI g-matrix  →  Critic + reward  (training only, never deployed)'),
]
for li, (ltr, clr, label, desc) in enumerate(layer_data):
    ly = y + SHH + 26 + li*30
    box(LX+5, ly, 14, 14, fill=clr)
    single(LX+5, ly+1, 14, 12, ltr, size=Pt(30), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    tf_l = tb(LX+22, ly, COL-27, 28)
    para(tf_l, label + ':  ' + desc, size=Pt(42), color=C_TEXT, first=True)

# Separator
box(LX+5, y+SHH+118, COL-10, 0.5, fill=C_BORDER)

# Model flow subheading
single(LX+5, y+SHH+122, COL-10, 14,
       'Model Flow', size=Pt(42), bold=True, color=C_PRIMARY)

# Flow boxes row 1: KPM → Encoder → z → broadcast
flow_items1 = [
    ('KPM\n[3×5]',  50, RGBColor(0xff,0xf3,0xcd)),
    ('Encoder\n(MLP→pool)', 65, RGBColor(0xd4,0xed,0xda)),
    ('z  [16]', 45, RGBColor(0xcc,0xe5,0xff)),
    ('→ broadcast\n  to all Actors', 80, C_CARD),
]
fx = LX+8
for txt, fw, fbg in flow_items1:
    if fbg != C_CARD:
        box(fx, y+SHH+140, fw, 22, fill=fbg, line=C_ACCENT, lw=Pt(0.3))
    single(fx+2, y+SHH+141, fw-4, 20, txt, size=Pt(28), bold=True,
           color=C_PRIMARY if fbg != C_CARD else C_ACCENT, align=PP_ALIGN.CENTER)
    fx += fw + (6 if fbg != C_CARD else 0)
# arrows
for ax in [LX+60, LX+133, LX+186]:
    single(ax, y+SHH+143, 8, 18, '→', size=Pt(36), color=C_ACCENT, align=PP_ALIGN.CENTER)

# Flow boxes row 2: local obs + z → SetActor → power
flow_items2 = [
    ('local obs\n[N_UE,i × 3]', 68, RGBColor(0xf8,0xd7,0xda)),
    ('+  z', 22, C_CARD),
    ('SetActor\n(perm-equivariant)', 80, RGBColor(0xe2,0xd9,0xf3)),
    ('power\n≤ P_max', 55, RGBColor(0xd1,0xec,0xf1)),
]
fx = LX+8
for txt, fw, fbg in flow_items2:
    if fbg != C_CARD:
        box(fx, y+SHH+170, fw, 22, fill=fbg, line=C_ACCENT, lw=Pt(0.3))
    single(fx+2, y+SHH+171, fw-4, 20, txt, size=Pt(28), bold=True,
           color=C_PRIMARY if fbg != C_CARD else C_ACCENT, align=PP_ALIGN.CENTER)
    fx += fw + 4
single(LX+194, y+SHH+173, 10, 18, '→', size=Pt(36), color=C_ACCENT, align=PP_ALIGN.CENTER)

# Key statement box
box(LX+5, y+SHH+200, COL-10, 22, fill=C_LIGHT, line=C_ACCENT, lw=Pt(0.4))
single(LX+9, y+SHH+203, COL-18, 18,
       'ONLY DIFFERENCE:  C-HASAC actor receives z   |   HASAC actor does NOT',
       size=Pt(40), bold=True, color=C_PRIMARY, align=PP_ALIGN.CENTER)

# HASAC theory
box(LX+5, y+SHH+228, COL-10, 0.5, fill=C_BORDER)
tf_h = tb(LX+5, y+SHH+233, COL-10, 112)
para(tf_h, 'HASAC: Sequential Soft Policy Decomposition (Liu et al., ICLR 2024 §3.3)',
     size=Pt(38), bold=True, color=C_PRIMARY, first=True)
para(tf_h, 'Each BS updates in random permutation order → converges to QRE',
     size=Pt(42), color=C_TEXT, space_before=4)
para(tf_h, '(Quantal Response Equilibrium = global MaxEnt optimum, not sub-optimal NE)',
     size=Pt(32), italic=True, color=C_MUTED, space_before=2)
para(tf_h, 'Empirical:  sequential HASAC −2.581  vs  simultaneous −5.184  (+2.6 PF-U ✓)',
     size=Pt(42), color=C_SUCCESS, space_before=4)

y += SHH + AH + SGAP

# ── EXPERIMENTAL SETUP  ───────────────────────────────────
EH = 220
section_hdr(LX, y, COL, 'Experimental Setup')
box(LX, y+SHH, COL, EH, fill=C_CARD, line=C_ACCENT, lw=Pt(0.4))

half = (COL-14) / 2
# Left half
single(LX+5, y+SHH+6, half, 13, 'Simulation (DeepMIMO)',
       size=Pt(38), bold=True, color=C_PRIMARY)
tf_sim = tb(LX+5, y+SHH+23, half, 88)
for line in ['•  N_BS = 3,  N_UE = 12',
             '•  Per-BS sum-power ≤ P_max',
             '•  Best-signal cell association',
             '•  Episode length T = 10 steps']:
    para(tf_sim, line, size=Pt(42), color=C_TEXT,
         first=(line == '•  N_BS = 3,  N_UE = 12'))

# Right half
single(LX+5+half+4, y+SHH+6, half, 13, 'Model & Training',
       size=Pt(38), bold=True, color=C_PRIMARY)
tf_mdl = tb(LX+5+half+4, y+SHH+23, half, 88)
for line in ['•  z_dim=16,  kpm_dim=5',
             '•  share_dim=63  (critic)',
             '•  Hidden=256,  SAC Twin-Q',
             '•  BC warm-start 1000 steps',
             '•  μ-bound=5  (tanh anti-collapse)']:
    para(tf_mdl, line, size=Pt(42), color=C_TEXT,
         first=(line == '•  z_dim=16,  kpm_dim=5'))

# PF formula
box(LX+5, y+SHH+116, COL-10, 0.5, fill=C_BORDER)
single(LX+5, y+SHH+120, COL-10, 12,
       'Evaluation Metric: PF-Utility (50 held-out episodes)',
       size=Pt(38), bold=True, color=C_PRIMARY, align=PP_ALIGN.CENTER)
single(LX+5, y+SHH+135, COL-10, 20,
       'U  =  Σu  log( R̄u + 10⁻⁶ )',
       size=Pt(52), bold=True, color=C_PRIMARY, align=PP_ALIGN.CENTER)

# Three benchmark cards
card_w = (COL-18) / 3
for ci, (lbl, val, clr) in enumerate([
    ('Floor (equal power)', '−5.332', C_MUTED),
    ('Best C-HASAC (400k)', '−1.162', C_SUCCESS),
    ('Ceiling (PF-WSR)', '+23.529', C_MUTED),
]):
    cx = LX+5 + ci*(card_w+4)
    box(cx, y+SHH+162, card_w, 44, fill=C_LIGHT, line=C_ACCENT, lw=Pt(0.3))
    single(cx+2, y+SHH+164, card_w-4, 22, val,
           size=Pt(46), bold=True, color=clr, align=PP_ALIGN.CENTER)
    single(cx+2, y+SHH+188, card_w-4, 16, lbl,
           size=Pt(26), italic=True, color=C_MUTED, align=PP_ALIGN.CENTER)

# ──────────────────────────────────────────────────────────
# RIGHT COLUMN
# ──────────────────────────────────────────────────────────
y = BY

# ── RESULTS  ──────────────────────────────────────────────
RH = 480
section_hdr(RX, y, COL, 'Results')
box(RX, y+SHH, COL, RH, fill=C_CARD, line=C_ACCENT, lw=Pt(0.4))

# Three metric cards
card_w3 = (COL-18) / 3
for ci, (val, lbl, clr) in enumerate([
    ('+1.42 PF-U', 'Gain: C-HASAC − HASAC', C_PRIMARY),
    ('+2.278', 'drop_shuffle (400k)', C_SUCCESS),
    ('−1.162', 'Best C-HASAC PF-U', C_SUCCESS),
]):
    cx = RX+5 + ci*(card_w3+4)
    box(cx, y+SHH+5, card_w3, 44, fill=C_LIGHT, line=C_ACCENT, lw=Pt(0.3))
    single(cx+2, y+SHH+7, card_w3-4, 22, val,
           size=Pt(46), bold=True, color=clr, align=PP_ALIGN.CENTER)
    single(cx+2, y+SHH+31, card_w3-4, 16, lbl,
           size=Pt(26), italic=True, color=C_MUTED, align=PP_ALIGN.CENTER)

# Results table
TH = y+SHH+56    # table header y
TRH = 14         # row height mm
COL_W = [118, 38, 38, 46, 68]   # column widths mm

# Header row
cx = RX+5
for cw, lbl in zip(COL_W, ['Method', 'PF-U (↑)', 'drop_zero', 'drop_shuffle', 'Note']):
    box(cx, TH, cw, TRH, fill=C_PRIMARY)
    single(cx+2, TH+2, cw-4, TRH-4, lbl, size=Pt(30), bold=True, color=C_WHITE)
    cx += cw

# Data rows
rows = [
    ('Equal Power',               '−5.332', '—',      '—',       'Floor',       C_BG,    C_MUTED,   False),
    ('HASAC (no z, seq 400k)',    '−2.581', '—',      '—',       'Sequential',  C_CARD,  C_TEXT,    False),
    ('C-HASAC geo_z (200k)',      '−2.237', '+0.932', '+1.429',  '—',           C_CARD,  C_TEXT,    False),
    ('C-HASAC + RSRP',           '−3.763', '+0.688', '+0.261',  'z ignored',   C_CARD,  C_WARNING, False),
    ('C-HASAC + Critic BC',      '−2.606', '−4.269', '−0.453',  'z harmful',   C_CARD,  C_WARNING, False),
    ('C-HASAC geo_z (400k) ★',   '−1.162', '+3.565', '+2.278',  '400k, BEST',  C_BESTBG,C_SUCCESS, True),
    ('PF-WSR oracle',            '+23.529','—',       '—',       'Full-CSI',    C_BG,    C_MUTED,   False),
]
for ri, (m, pu, dz, ds, note, rbg, rfg, rb) in enumerate(rows):
    ry = TH + TRH + ri*TRH
    cx = RX+5
    for cw, val in zip(COL_W, [m, pu, dz, ds, note]):
        box(cx, ry, cw, TRH, fill=rbg, line=C_BORDER, lw=Pt(0.2))
        single(cx+2, ry+1, cw-4, TRH-2, val,
               size=Pt(28) if not rb else Pt(30), bold=rb, color=rfg)
        cx += cw

# Z-ablation chart
CY = TH + TRH*(1+len(rows)) + 8
single(RX+5, CY, COL-10, 12,
       'Z-Ablation:  drop_shuffle by Configuration',
       size=Pt(38), bold=True, color=C_PRIMARY, align=PP_ALIGN.CENTER)

CHART_BASE = CY+14         # top of chart area mm
ZERO_OFF   = 42            # mm from CHART_BASE to zero-line
SCALE      = 32            # mm per 1.0 PF-U unit
BW         = 44            # bar width mm
BGAP       = 18

# zero axis
box(RX+18, CHART_BASE+ZERO_OFF, COL-23, 0.6, fill=RGBColor(0xaa,0xaa,0xaa))
for yoff, lbl in [(-ZERO_OFF, ''), (-SCALE, '+1'), (-2*SCALE+4, '+2')]:
    single(RX+5, CHART_BASE+ZERO_OFF+yoff, 12, 8, lbl,
           size=Pt(26), color=C_MUTED, align=PP_ALIGN.RIGHT)

chart_bars = [
    ('+1.43', 1.429, C_ACCENT,   'geo_z\n200k'),
    ('+2.28', 2.278, C_PRIMARY,  'geo_z\n400k ★'),
    ('+0.26', 0.261, C_ORANGE,   '+RSRP'),
    ('−0.45', -0.453,C_WARNING,  '+Critic\nBC'),
]
bx = RX+30
for val_lbl, val, clr, name in chart_bars:
    bh = abs(val) * SCALE
    by = (CHART_BASE + ZERO_OFF - bh) if val >= 0 else (CHART_BASE + ZERO_OFF)
    box(bx, by, BW, bh, fill=clr)
    # value label
    vly = by - 10 if val >= 0 else by + bh + 2
    single(bx, vly, BW, 9, val_lbl, size=Pt(28), bold=True, color=clr, align=PP_ALIGN.CENTER)
    # name label (below zero line always)
    name_y = CHART_BASE + ZERO_OFF + (bh if val < 0 else 0) + 12
    single(bx-4, name_y, BW+8, 16, name, size=Pt(26), color=C_TEXT, align=PP_ALIGN.CENTER)
    bx += BW + BGAP

y += SHH + RH + SGAP

# ── ANALYSIS & DISCUSSION  ────────────────────────────────
ANAH = 300
section_hdr(RX, y, COL, 'Analysis & Discussion')
box(RX, y+SHH, COL, ANAH, fill=C_CARD, line=C_ACCENT, lw=Pt(0.4))

half_r = (COL-14) / 2
# Left: Why z works
single(RX+5, y+SHH+6, half_r, 14, 'Why z works (geo_z)',
       size=Pt(38), bold=True, color=C_SUCCESS)
tf_w = tb(RX+5, y+SHH+24, half_r, 140)
lines_w = [
    ('• BC warm-start unlocks z:', True, C_TEXT),
    ('  pure RL → drop_shuffle ≈ 0', False, C_TEXT),
    ('  +BC 1000 steps → +2.278 ✓', False, C_SUCCESS),
    ('• KPM with BS distances →', True, C_TEXT),
    ('  z encodes spatial topology', False, C_TEXT),
    ('• drop_shuffle > drop_zero →', True, C_TEXT),
    ('  z content used, not offset', False, C_TEXT),
    ('• 400k > 200k: +1.43 → +2.28', False, C_ACCENT),
]
for i, (t, bo, cl) in enumerate(lines_w):
    para(tf_w, t, size=Pt(42), bold=bo, color=cl, first=(i==0), space_before=(4 if bo else 0))

# Right: Why ablations fail
single(RX+5+half_r+4, y+SHH+6, half_r, 14, 'Why ablations fail',
       size=Pt(38), bold=True, color=C_WARNING)
tf_f = tb(RX+5+half_r+4, y+SHH+24, half_r, 140)
lines_f = [
    ('• +RSRP: actor gets neighbor', True, C_TEXT),
    ('  channel gains → z ignored', False, C_TEXT),
    ('  +1.43 → +0.26 (unused)', False, C_WARNING),
    ('• +Critic BC: encoder learns', True, C_TEXT),
    ('  to satisfy pre-trained Q,', False, C_TEXT),
    ('  not actor → z misleads policy', False, C_TEXT),
    ('  drop_shuffle = −0.45 ✗', True, C_WARNING),
]
for i, (t, bo, cl) in enumerate(lines_f):
    para(tf_f, t, size=Pt(42), bold=bo, color=cl, first=(i==0), space_before=(4 if bo else 0))

box(RX+5, y+SHH+172, COL-10, 0.5, fill=C_BORDER)
single(RX+5, y+SHH+177, COL-10, 13,
       'Open Problem: SAC Q-Overestimation',
       size=Pt(38), bold=True, color=C_PRIMARY)
tf_op = tb(RX+5, y+SHH+194, COL-10, 100)
para(tf_op, '• Peaks at step ~40–80k then oscillates — max operator amplifies estimation error.',
     size=Pt(42), color=C_TEXT, first=True)
para(tf_op, '• Mitigated: μ-bound=5 + logpf reward + BC warm-start.',
     size=Pt(42), color=C_TEXT, space_before=4)
para(tf_op, '• Future work: Conservative Q-Learning (CQL) penalty.',
     size=Pt(42), italic=True, color=C_ACCENT, space_before=4)

y += SHH + ANAH + SGAP

# ── CONCLUSION  ───────────────────────────────────────────
CON_H = 255
section_hdr(RX, y, COL, 'Conclusion')
box(RX, y+SHH, COL, CON_H, fill=C_CARD, line=C_ACCENT, lw=Pt(0.4))

tf_c = tb(RX+5, y+SHH+8, COL-10, CON_H-16)
concl = [
    ('✅  C-HASAC > HASAC by +1.42 PF-U  (−1.162 vs −2.581, 400k steps).', True, C_TEXT),
    ('    The only difference: actor receives learned latent context z.', False, C_TEXT),
    ('✅  z is genuinely used: drop_shuffle = +2.278.', True, C_TEXT),
    ('    Mismatched z hurts more than zeroing — content matters.', False, C_TEXT),
    ('✅  BC warm-start is the key catalyst (pure RL → drop_shuffle ≈ 0).', True, C_TEXT),
    ('    1000 expert imitation steps opens the z-usage switch.', False, C_TEXT),
    ('✅  O-RAN deployable: RIC xApp → E2 interface → gNBs.', True, C_TEXT),
    ('    No direct inter-BS CSI exchange required at runtime.', False, C_TEXT),
    ('⚠️  Gap to PF-WSR ceiling (+23.529) remains; future: CQL for Q stability.', False, C_MUTED),
]
for i, (t, bo, cl) in enumerate(concl):
    para(tf_c, t, size=Pt(42), bold=bo, color=cl,
         first=(i==0), space_before=(8 if bo and i>0 else 0))

# ══════════════════════════════════════════════════════════
prs.save(os.path.join(HERE, 'poster.pptx'))
print('Saved: poster/poster.pptx')
